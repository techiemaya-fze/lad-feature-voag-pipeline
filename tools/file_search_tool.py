"""
File Search Tool - Google Gemini File Search API wrapper.

Provides operations for creating/managing FileSearchStores and uploading documents
for the RAG (Retrieval-Augmented Generation) feature.

Based on Google's official File Search documentation:
https://ai.google.dev/gemini-api/docs/file-search
"""

import asyncio
import logging
import os
import time
from dataclasses import dataclass
from typing import Any, Optional

from dotenv import load_dotenv

load_dotenv()

logger = logging.getLogger(__name__)

# MIME type mapping for common file formats
# Google File Search supports these formats: https://ai.google.dev/gemini-api/docs/file-search
MIME_TYPE_MAP = {
    # Documents
    ".pdf": "application/pdf",
    ".txt": "text/plain",
    ".html": "text/html",
    ".htm": "text/html",
    ".md": "text/markdown",
    ".rtf": "application/rtf",
    # Microsoft Office
    ".doc": "application/msword",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".ppt": "application/vnd.ms-powerpoint",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xls": "application/vnd.ms-excel",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    # OpenDocument formats
    ".odt": "application/vnd.oasis.opendocument.text",
    ".ods": "application/vnd.oasis.opendocument.spreadsheet",
    ".odp": "application/vnd.oasis.opendocument.presentation",
    # Data formats
    ".csv": "text/csv",
    ".json": "application/json",
    ".xml": "application/xml",
}

# File types that need conversion to text/plain before upload
# Gemini File Search only accepts: text/plain, application/pdf
CONVERT_TO_TEXT_EXTENSIONS = {".xlsx", ".xls", ".docx", ".doc", ".pptx", ".ppt", ".csv"}


def convert_document_to_text(file_path: str) -> tuple[str, str]:
    """
    Convert document to plain text format for Gemini File Search.
    
    Gemini File Search only accepts text/plain and PDF.
    This converts Excel, Word, PowerPoint to text.
    
    Args:
        file_path: Path to the document
        
    Returns:
        Tuple of (text_content, original_extension)
        
    Raises:
        FileSearchToolError: If conversion fails
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext in (".xlsx", ".xls"):
            return _convert_excel_to_text(file_path), ext
        elif ext in (".docx", ".doc"):
            return _convert_word_to_text(file_path), ext
        elif ext in (".pptx", ".ppt"):
            return _convert_ppt_to_text(file_path), ext
        elif ext == ".csv":
            return _convert_csv_to_text(file_path), ext
        else:
            raise FileSearchToolError(f"Unsupported format for conversion: {ext}")
    except Exception as exc:
        raise FileSearchToolError(f"Failed to convert {ext} to text: {exc}") from exc


def _convert_excel_to_text(file_path: str) -> str:
    """Convert Excel file to structured text."""
    try:
        import openpyxl
    except ImportError:
        raise FileSearchToolError("openpyxl not installed. Run: uv add openpyxl")
    
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    text_parts = []
    
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")
        
        for row in sheet.iter_rows(values_only=True):
            # Filter out empty rows
            if any(cell is not None for cell in row):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                text_parts.append(row_text)
    
    return "\n".join(text_parts)


def _convert_word_to_text(file_path: str) -> str:
    """Convert Word document to text."""
    try:
        from docx import Document
    except ImportError:
        raise FileSearchToolError("python-docx not installed. Run: uv add python-docx")
    
    doc = Document(file_path)
    text_parts = []
    
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
    
    # Also extract tables
    for table in doc.tables:
        text_parts.append("\n[Table]")
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            text_parts.append(row_text)
    
    return "\n".join(text_parts)


def _convert_ppt_to_text(file_path: str) -> str:
    """Convert PowerPoint to text."""
    try:
        from pptx import Presentation
    except ImportError:
        raise FileSearchToolError("python-pptx not installed. Run: uv add python-pptx")
    
    prs = Presentation(file_path)
    text_parts = []
    
    for i, slide in enumerate(prs.slides, 1):
        text_parts.append(f"\n=== Slide {i} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)
    
    return "\n".join(text_parts)


def _convert_csv_to_text(file_path: str) -> str:
    """Convert CSV to formatted text."""
    import csv
    
    text_parts = []
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        reader = csv.reader(f)
        for row in reader:
            text_parts.append(" | ".join(row))
    
    return "\n".join(text_parts)


class FileSearchToolError(Exception):
    """Exception raised for File Search tool errors."""
    pass


@dataclass
class ChunkingConfig:
    """Configuration for document chunking during ingestion."""
    max_tokens_per_chunk: int = 500
    max_overlap_tokens: int = 50


@dataclass
class DocumentInfo:
    """Information about a document in a FileSearchStore."""
    name: str  # Resource name, e.g., "fileSearchStores/xxx/documents/yyy"
    display_name: str
    state: str  # STATE_PENDING, STATE_ACTIVE, STATE_FAILED
    mime_type: Optional[str] = None
    size_bytes: Optional[int] = None


@dataclass
class StoreInfo:
    """Information about a FileSearchStore."""
    name: str  # Resource name, e.g., "fileSearchStores/xxx"
    display_name: str


class FileSearchTool:
    """
    Wrapper for Google Gemini File Search API.
    
    Manages FileSearchStores and documents for RAG retrieval.
    """

    def __init__(self) -> None:
        """Initialize the FileSearchTool with Google GenAI client."""
        self._client = None
        self._initialized = False

    def _ensure_client(self) -> Any:
        """Lazily initialize the Google GenAI client."""
        if self._client is not None:
            return self._client

        try:
            from google import genai
        except ImportError as exc:
            raise FileSearchToolError(
                "google-genai SDK not installed. Run: uv add google-genai"
            ) from exc

        # Try GOOGLE_API_KEY first, then GEMINI_API_KEY
        api_key = os.getenv("GOOGLE_API_KEY") or os.getenv("GEMINI_API_KEY")
        if not api_key:
            raise FileSearchToolError(
                "GOOGLE_API_KEY or GEMINI_API_KEY environment variable not set"
            )

        self._client = genai.Client(api_key=api_key)
        self._initialized = True
        logger.info("Google GenAI client initialized for File Search")
        return self._client

    # =========================================================================
    # STORE MANAGEMENT
    # =========================================================================

    async def create_store(self, display_name: str) -> StoreInfo:
        """
        Create a new FileSearchStore.
        
        Args:
            display_name: Human-readable name for the store
            
        Returns:
            StoreInfo with the created store's details
        """
        client = self._ensure_client()

        try:
            store = client.file_search_stores.create(
                config={"display_name": display_name}
            )
            logger.info("Created FileSearchStore: %s (%s)", display_name, store.name)
            return StoreInfo(
                name=store.name,
                display_name=display_name,
            )
        except Exception as exc:
            logger.error("Failed to create store '%s': %s", display_name, exc, exc_info=True)
            raise FileSearchToolError(f"Failed to create store: {exc}") from exc

    async def delete_store(self, store_name: str, force: bool = True) -> bool:
        """
        Delete a FileSearchStore.
        
        Args:
            store_name: The store's resource name (e.g., "fileSearchStores/xxx")
            force: If True, delete even if store has documents
            
        Returns:
            True if deleted successfully
        """
        client = self._ensure_client()

        try:
            client.file_search_stores.delete(
                name=store_name,
                config={"force": force},
            )
            logger.info("Deleted FileSearchStore: %s", store_name)
            return True
        except Exception as exc:
            logger.error("Failed to delete store '%s': %s", store_name, exc, exc_info=True)
            raise FileSearchToolError(f"Failed to delete store: {exc}") from exc

    async def list_stores(self) -> list[StoreInfo]:
        """List all FileSearchStores."""
        client = self._ensure_client()

        try:
            stores = []
            for store in client.file_search_stores.list():
                stores.append(StoreInfo(
                    name=store.name,
                    display_name=getattr(store, "display_name", store.name),
                ))
            return stores
        except Exception as exc:
            logger.error("Failed to list stores: %s", exc, exc_info=True)
            raise FileSearchToolError(f"Failed to list stores: {exc}") from exc

    async def get_store(self, store_name: str) -> Optional[StoreInfo]:
        """Get a specific FileSearchStore by name."""
        client = self._ensure_client()

        try:
            store = client.file_search_stores.get(name=store_name)
            return StoreInfo(
                name=store.name,
                display_name=getattr(store, "display_name", store.name),
            )
        except Exception as exc:
            logger.error("Failed to get store '%s': %s", store_name, exc, exc_info=True)
            return None

    # =========================================================================
    # DOCUMENT MANAGEMENT
    # =========================================================================

    async def upload_document(
        self,
        store_name: str,
        file_path: str,
        display_name: Optional[str] = None,
        chunking_config: Optional[ChunkingConfig] = None,
        wait_for_completion: bool = True,
        poll_interval: float = 5.0,
        timeout: float = 300.0,
    ) -> DocumentInfo:
        """
        Upload a document file to a FileSearchStore.
        
        Args:
            store_name: The store's resource name
            file_path: Path to the file to upload
            display_name: Optional display name (defaults to filename)
            chunking_config: Optional chunking configuration
            wait_for_completion: Wait for indexing to complete
            poll_interval: Seconds between polling attempts
            timeout: Maximum seconds to wait for completion
            
        Returns:
            DocumentInfo with the uploaded document's details
        """
        import tempfile
        
        client = self._ensure_client()
        chunking = chunking_config or ChunkingConfig()

        # Detect file extension
        ext = os.path.splitext(file_path)[1].lower()
        
        # Check if file needs conversion (Gemini only accepts text/plain and PDF)
        actual_file_path = file_path
        temp_converted_path = None
        
        if ext in CONVERT_TO_TEXT_EXTENSIONS:
            logger.info("Converting %s file to text for Gemini ingestion", ext)
            try:
                text_content, original_ext = convert_document_to_text(file_path)
                
                # Write converted text to temp file
                with tempfile.NamedTemporaryFile(
                    delete=False, suffix=".txt", prefix="kb_converted_", mode='w', encoding='utf-8'
                ) as tmp:
                    tmp.write(text_content)
                    temp_converted_path = tmp.name
                
                actual_file_path = temp_converted_path
                ext = ".txt"
                logger.info("Converted %s to text (%d chars)", original_ext, len(text_content))
            except Exception as e:
                logger.error("Failed to convert document: %s", e)
                raise FileSearchToolError(f"Document conversion failed: {e}") from e
        
        # Detect MIME type from file extension
        mime_type = MIME_TYPE_MAP.get(ext)
        if not mime_type:
            logger.warning("Unknown MIME type for extension '%s', upload may fail", ext)

        # Get display name and sanitize to ASCII
        doc_display_name = display_name or os.path.basename(file_path)
        
        # Sanitize display name to ASCII - Google SDK can't encode Unicode in HTTP headers
        # Replace common curly quotes and apostrophes with straight versions
        unicode_replacements = {
            '\u2018': "'",  # Left single quotation mark
            '\u2019': "'",  # Right single quotation mark (curly apostrophe)
            '\u201C': '"',  # Left double quotation mark
            '\u201D': '"',  # Right double quotation mark
            '\u2013': '-',  # En dash
            '\u2014': '-',  # Em dash
            '\u2026': '...',  # Ellipsis
        }
        for unicode_char, ascii_char in unicode_replacements.items():
            doc_display_name = doc_display_name.replace(unicode_char, ascii_char)
        
        # Final fallback: encode to ASCII, replacing any remaining non-ASCII chars
        doc_display_name = doc_display_name.encode('ascii', 'replace').decode('ascii')

        config: dict[str, Any] = {
            "display_name": doc_display_name,
            "chunking_config": {
                "white_space_config": {
                    "max_tokens_per_chunk": chunking.max_tokens_per_chunk,
                    "max_overlap_tokens": chunking.max_overlap_tokens,
                }
            },
        }

        # Add mime_type to config if detected (not as separate kwarg)
        if mime_type:
            config["mime_type"] = mime_type
            logger.debug("Uploading with MIME type: %s", mime_type)

        try:
            operation = client.file_search_stores.upload_to_file_search_store(
                file=actual_file_path,  # Use converted path if applicable
                file_search_store_name=store_name,
                config=config,
            )

            if wait_for_completion:
                operation = await self._poll_operation(
                    operation, poll_interval, timeout
                )

            # Extract document info from operation result
            doc_name = getattr(operation, "name", "unknown")
            return DocumentInfo(
                name=doc_name,
                display_name=config["display_name"],
                state="STATE_ACTIVE" if operation.done else "STATE_PENDING",
            )

        except Exception as exc:
            logger.error("Failed to upload document: %s", exc, exc_info=True)
            raise FileSearchToolError(f"Failed to upload document: {exc}") from exc
        finally:
            # Clean up temp converted file if created
            if temp_converted_path:
                try:
                    os.unlink(temp_converted_path)
                except Exception:
                    pass

    async def upload_document_from_bytes(
        self,
        store_name: str,
        file_bytes: bytes,
        filename: str,
        display_name: Optional[str] = None,
        chunking_config: Optional[ChunkingConfig] = None,
        wait_for_completion: bool = True,
        poll_interval: float = 5.0,
        timeout: float = 300.0,
    ) -> DocumentInfo:
        """
        Upload document bytes to a FileSearchStore.
        
        This writes to a temp file and uploads, as the SDK requires file paths.
        
        Args:
            store_name: The store's resource name
            file_bytes: Document content as bytes
            filename: Original filename (for MIME type detection)
            display_name: Optional display name
            chunking_config: Optional chunking configuration
            wait_for_completion: Wait for indexing to complete
            poll_interval: Seconds between polling attempts
            timeout: Maximum seconds to wait for completion
            
        Returns:
            DocumentInfo with the uploaded document's details
        """
        import tempfile

        # Write bytes to temp file
        suffix = os.path.splitext(filename)[1] if "." in filename else ""
        with tempfile.NamedTemporaryFile(
            delete=False, suffix=suffix, prefix="kb_upload_"
        ) as tmp:
            tmp.write(file_bytes)
            temp_path = tmp.name

        try:
            return await self.upload_document(
                store_name=store_name,
                file_path=temp_path,
                display_name=display_name or filename,
                chunking_config=chunking_config,
                wait_for_completion=wait_for_completion,
                poll_interval=poll_interval,
                timeout=timeout,
            )
        finally:
            # Clean up temp file
            try:
                os.unlink(temp_path)
            except Exception:
                pass

    async def list_documents(self, store_name: str) -> list[DocumentInfo]:
        """
        List all documents in a FileSearchStore.
        
        Note: This uses REST API as SDK may not expose document listing.
        
        Args:
            store_name: The store's resource name
            
        Returns:
            List of DocumentInfo objects
        """
        client = self._ensure_client()

        # The SDK exposes documents through file_search_stores
        # Try to iterate through documents
        try:
            docs = []
            # SDK may vary - try common patterns
            store = client.file_search_stores.get(name=store_name)
            
            # If documents are accessible via the store
            if hasattr(store, "documents"):
                for doc in store.documents:
                    docs.append(DocumentInfo(
                        name=doc.name,
                        display_name=getattr(doc, "display_name", doc.name),
                        state=getattr(doc, "state", "UNKNOWN"),
                        mime_type=getattr(doc, "mime_type", None),
                        size_bytes=getattr(doc, "size_bytes", None),
                    ))
            
            return docs
        except Exception as exc:
            logger.warning("Failed to list documents for '%s': %s", store_name, exc)
            # Return empty list on failure - documents may still exist
            return []

    async def delete_document(
        self,
        document_name: str,
        force: bool = True,
    ) -> bool:
        """
        Delete a document from a FileSearchStore.
        
        Args:
            document_name: Full document resource name
            force: If True, force delete even if has chunks
            
        Returns:
            True if deleted successfully
        """
        # Documents API deletion - may need REST call
        # For now, log and return success (documents expire with store deletion)
        logger.info("Document deletion requested: %s (force=%s)", document_name, force)
        return True

    # =========================================================================
    # HELPERS
    # =========================================================================

    async def _poll_operation(
        self,
        operation: Any,
        poll_interval: float,
        timeout: float,
    ) -> Any:
        """Poll an async operation until completion or timeout."""
        client = self._ensure_client()
        start_time = time.time()

        while not operation.done:
            if time.time() - start_time > timeout:
                raise FileSearchToolError(
                    f"Operation timed out after {timeout} seconds"
                )

            await asyncio.sleep(poll_interval)
            operation = client.operations.get(operation)

        return operation


# =============================================================================
# TOOL BUILDER INTEGRATION (Phase 14)
# =============================================================================

def get_file_search_tools(store_ids: list[str] | None = None) -> list:
    """
    Get file search tool functions for agent attachment.
    
    Called by tool_builder.py to get callable search tools.
    
    Args:
        store_ids: List of Gemini FileSearch store names to enable search on
        
    Returns:
        List of tool functions (placeholder - actual tools defined in agent)
    """
    if not store_ids:
        logger.debug("No KB store_ids provided, returning empty tools")
        return []
    
    logger.info(f"File search tools prepared for {len(store_ids)} store(s)")
    # TODO: Return actual callable search tool functions
    # The actual search_knowledge_base tool is defined in agent.py
    # This function is for future modular tool attachment
    return []

"""
Document Converter - Convert Office files to PDF or structured text.

Provides conversion pipeline for uploading Office documents to Google File Search
when the API rejects native Office formats.

Conversion strategies (in order of preference):
1. LibreOffice (headless) - Best quality, preserves all formatting
2. Structured text extraction - Fallback when no converter available

Install LibreOffice for best results:
- Ubuntu/Debian: apt install libreoffice
- Windows: choco install libreoffice or download from libreoffice.org
- macOS: brew install --cask libreoffice
"""

import asyncio
import io
import logging
import os
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Optional, Tuple, List

logger = logging.getLogger(__name__)


@dataclass
class ConversionResult:
    """Result of a document conversion."""
    output_path: str
    output_mime_type: str
    original_path: str
    original_mime_type: str
    conversion_method: str  # 'libreoffice', 'textextract', 'none'
    success: bool
    error_message: Optional[str] = None


class DocumentConverter:
    """
    Convert Office documents to formats supported by Google File Search.
    
    Primary strategy: Convert to PDF using LibreOffice (preserves structure).
    Fallback: Extract structured text (for when no converter is available).
    """
    
    # MIME types for Office formats
    OFFICE_MIME_TYPES = {
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".doc": "application/msword",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".ppt": "application/vnd.ms-powerpoint",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xls": "application/vnd.ms-excel",
    }
    
    # Formats that need conversion
    NEEDS_CONVERSION = {".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"}
    
    def __init__(self, temp_dir: Optional[str] = None):
        """Initialize converter with optional temp directory for outputs."""
        self.temp_dir = temp_dir or tempfile.gettempdir()
        self._libreoffice_path = self._find_libreoffice()
        
        if self._libreoffice_path:
            logger.info("DocumentConverter: LibreOffice found at %s", self._libreoffice_path)
        else:
            logger.warning("DocumentConverter: LibreOffice not found - will use text extraction fallback")
    
    def _find_libreoffice(self) -> Optional[str]:
        """Find LibreOffice executable path."""
        # Common paths on different platforms
        paths_to_try = [
            # Windows
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            # Linux
            "/usr/bin/libreoffice",
            "/usr/bin/soffice",
            "/opt/libreoffice/program/soffice",
            # macOS
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        ]
        
        for path in paths_to_try:
            if os.path.exists(path):
                return path
        
        # Try to find via PATH
        try:
            result = subprocess.run(
                ["which", "libreoffice"] if os.name != "nt" else ["where", "soffice"],
                capture_output=True,
                text=True,
                timeout=5
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip().split("\n")[0]
        except Exception:
            pass
        
        return None
    
    def needs_conversion(self, file_path: str) -> bool:
        """Check if a file needs conversion for File Search upload."""
        ext = Path(file_path).suffix.lower()
        return ext in self.NEEDS_CONVERSION
    
    def can_convert_to_pdf(self) -> bool:
        """Check if PDF conversion is available (LibreOffice installed)."""
        return self._libreoffice_path is not None
    
    async def convert(
        self,
        file_path: str,
        prefer_pdf: bool = True,
        output_dir: Optional[str] = None,
    ) -> ConversionResult:
        """
        Convert an Office document to a format supported by File Search.
        
        Args:
            file_path: Path to the Office document
            prefer_pdf: If True, prefer PDF output (requires LibreOffice)
            output_dir: Directory for output file (default: temp directory)
            
        Returns:
            ConversionResult with output path and metadata
        """
        path = Path(file_path)
        ext = path.suffix.lower()
        
        if ext not in self.NEEDS_CONVERSION:
            # File doesn't need conversion
            return ConversionResult(
                output_path=file_path,
                output_mime_type=self._get_mime_type(file_path),
                original_path=file_path,
                original_mime_type=self._get_mime_type(file_path),
                conversion_method="none",
                success=True,
            )
        
        output_dir = output_dir or self.temp_dir
        original_mime = self.OFFICE_MIME_TYPES.get(ext, "application/octet-stream")
        
        # Try LibreOffice PDF conversion first if preferred and available
        if prefer_pdf and self._libreoffice_path:
            result = await self._convert_with_libreoffice(file_path, output_dir)
            if result.success:
                return result
            logger.warning("LibreOffice conversion failed, falling back to text extraction")
        
        # Fallback: Extract structured text
        return await self._extract_text(file_path, output_dir)
    
    async def _convert_with_libreoffice(
        self,
        file_path: str,
        output_dir: str,
    ) -> ConversionResult:
        """Convert to PDF using LibreOffice headless mode."""
        path = Path(file_path)
        ext = path.suffix.lower()
        original_mime = self.OFFICE_MIME_TYPES.get(ext, "application/octet-stream")
        
        try:
            # LibreOffice command for PDF conversion
            cmd = [
                self._libreoffice_path,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", output_dir,
                str(path.resolve()),
            ]
            
            logger.info("Converting %s to PDF with LibreOffice...", path.name)
            
            # Run in executor to avoid blocking
            loop = asyncio.get_event_loop()
            result = await loop.run_in_executor(
                None,
                lambda: subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=120,  # 2 minute timeout
                )
            )
            
            if result.returncode != 0:
                error = result.stderr or result.stdout or "Unknown error"
                logger.error("LibreOffice conversion failed: %s", error)
                return ConversionResult(
                    output_path="",
                    output_mime_type="",
                    original_path=file_path,
                    original_mime_type=original_mime,
                    conversion_method="libreoffice",
                    success=False,
                    error_message=error,
                )
            
            # Output file has same name with .pdf extension
            output_path = os.path.join(output_dir, path.stem + ".pdf")
            
            if not os.path.exists(output_path):
                return ConversionResult(
                    output_path="",
                    output_mime_type="",
                    original_path=file_path,
                    original_mime_type=original_mime,
                    conversion_method="libreoffice",
                    success=False,
                    error_message="Output PDF not created",
                )
            
            logger.info("Successfully converted to PDF: %s", output_path)
            
            return ConversionResult(
                output_path=output_path,
                output_mime_type="application/pdf",
                original_path=file_path,
                original_mime_type=original_mime,
                conversion_method="libreoffice",
                success=True,
            )
            
        except subprocess.TimeoutExpired:
            return ConversionResult(
                output_path="",
                output_mime_type="",
                original_path=file_path,
                original_mime_type=original_mime,
                conversion_method="libreoffice",
                success=False,
                error_message="Conversion timed out",
            )
        except Exception as e:
            logger.error("LibreOffice conversion error: %s", e)
            return ConversionResult(
                output_path="",
                output_mime_type="",
                original_path=file_path,
                original_mime_type=original_mime,
                conversion_method="libreoffice",
                success=False,
                error_message=str(e),
            )
    
    async def _extract_text(
        self,
        file_path: str,
        output_dir: str,
    ) -> ConversionResult:
        """Extract structured text from Office documents."""
        path = Path(file_path)
        ext = path.suffix.lower()
        original_mime = self.OFFICE_MIME_TYPES.get(ext, "application/octet-stream")
        
        try:
            if ext in {".docx", ".doc"}:
                text = await self._extract_docx(file_path)
                output_ext = ".txt"
                output_mime = "text/plain"
            elif ext in {".pptx", ".ppt"}:
                text = await self._extract_pptx(file_path)
                output_ext = ".txt"  
                output_mime = "text/plain"
            elif ext in {".xlsx", ".xls"}:
                text = await self._extract_xlsx(file_path)
                output_ext = ".csv"  # CSV format for spreadsheets
                output_mime = "text/csv"
            else:
                return ConversionResult(
                    output_path="",
                    output_mime_type="",
                    original_path=file_path,
                    original_mime_type=original_mime,
                    conversion_method="textextract",
                    success=False,
                    error_message=f"Unsupported format: {ext}",
                )
            
            # Write extracted text to file
            output_path = os.path.join(output_dir, path.stem + output_ext)
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(text)
            
            logger.info("Extracted text to: %s (%d chars)", output_path, len(text))
            
            return ConversionResult(
                output_path=output_path,
                output_mime_type=output_mime,
                original_path=file_path,
                original_mime_type=original_mime,
                conversion_method="textextract",
                success=True,
            )
            
        except Exception as e:
            logger.error("Text extraction error: %s", e, exc_info=True)
            return ConversionResult(
                output_path="",
                output_mime_type="",
                original_path=file_path,
                original_mime_type=original_mime,
                conversion_method="textextract",
                success=False,
                error_message=str(e),
            )
    
    async def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX files with structure preservation."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx required: pip install python-docx")
        
        doc = Document(file_path)
        parts = []
        
        # Document title/heading
        parts.append(f"# {Path(file_path).stem}")
        parts.append("")
        
        # Paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                # Check for heading styles
                if para.style and "Heading" in para.style.name:
                    level = 1
                    if "1" in para.style.name:
                        level = 1
                    elif "2" in para.style.name:
                        level = 2
                    elif "3" in para.style.name:
                        level = 3
                    parts.append(f"{'#' * level} {para.text}")
                else:
                    parts.append(para.text)
                parts.append("")
        
        # Tables
        for table_idx, table in enumerate(doc.tables):
            parts.append(f"\n## Table {table_idx + 1}")
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                parts.append(row_text)
            parts.append("")
        
        return "\n".join(parts)
    
    async def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX files with slide structure."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx required: pip install python-pptx")
        
        prs = Presentation(file_path)
        parts = []
        
        # Presentation title
        parts.append(f"# {Path(file_path).stem}")
        parts.append("")
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            parts.append(f"## Slide {slide_num}")
            parts.append("")
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
                    parts.append("")
                
                # Tables in slides
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        parts.append(row_text)
                    parts.append("")
            
            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"**Speaker Notes:** {notes}")
                    parts.append("")
            
            parts.append("---")
            parts.append("")
        
        return "\n".join(parts)
    
    async def _extract_xlsx(self, file_path: str) -> str:
        """Extract data from XLSX files as CSV format."""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl required: pip install openpyxl")
        
        import csv
        from io import StringIO
        
        wb = openpyxl.load_workbook(file_path, data_only=True)
        output = StringIO()
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Add sheet header
            output.write(f"\n# Sheet: {sheet_name}\n")
            
            writer = csv.writer(output)
            
            for row in ws.iter_rows(values_only=True, max_row=2000):
                # Convert None to empty string
                row_data = ["" if cell is None else str(cell) for cell in row]
                # Skip completely empty rows
                if any(cell for cell in row_data):
                    writer.writerow(row_data)
        
        return output.getvalue()
    
    def _get_mime_type(self, file_path: str) -> str:
        """Get MIME type for a file."""
        ext = Path(file_path).suffix.lower()
        mime_map = {
            ".pdf": "application/pdf",
            ".txt": "text/plain",
            ".csv": "text/csv",
            ".md": "text/markdown",
            ".html": "text/html",
            **self.OFFICE_MIME_TYPES,
        }
        return mime_map.get(ext, "application/octet-stream")
